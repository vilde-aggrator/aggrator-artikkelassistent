import anthropic
import base64
import json
import os
import re
import smtplib
from datetime import datetime
from email.header import Header
from email.mime.image import MIMEImage
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from io import BytesIO
from pathlib import Path
from typing import Optional

BASE_DIR = Path(__file__).parent
GUIDELINES_DIR = BASE_DIR / "_guidelines"

client = anthropic.Anthropic()


def _load_guidelines() -> str:
    files = ["aggrator-stil.md", "tone-of-voice.md", "mal-artikkel.md"]
    parts = []
    for f in files:
        p = GUIDELINES_DIR / f
        if p.exists():
            parts.append(f"### {f}\n\n{p.read_text()}")
    return "\n\n---\n\n".join(parts)


GUIDELINES = _load_guidelines()


def _process_files(files_data: list) -> tuple[list, str]:
    """Returns (pdf_document_blocks, extracted_text) for uploaded files."""
    doc_blocks = []
    text_parts = []

    for f in files_data:
        name = f["name"]
        fbytes = f["bytes"]

        if name.lower().endswith(".pdf"):
            doc_blocks.append({
                "type": "document",
                "source": {
                    "type": "base64",
                    "media_type": "application/pdf",
                    "data": base64.standard_b64encode(fbytes).decode("utf-8"),
                },
                "title": name,
            })
        elif name.lower().endswith(".docx"):
            try:
                from docx import Document
                doc = Document(BytesIO(fbytes))
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                text_parts.append(f"[{name}]\n{text}")
            except Exception:
                pass
        elif name.lower().endswith(".pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(BytesIO(fbytes))
                slides = []
                for i, slide in enumerate(prs.slides):
                    lines = [
                        shape.text.strip()
                        for shape in slide.shapes
                        if hasattr(shape, "text") and shape.text.strip()
                    ]
                    if lines:
                        slides.append(f"[Slide {i + 1}]\n" + "\n".join(lines))
                text_parts.append(f"[{name}]\n" + "\n\n".join(slides))
            except Exception:
                pass
        elif name.lower().endswith(".txt"):
            text_parts.append(f"[{name}]\n{fbytes.decode('utf-8', errors='replace')}")

    return doc_blocks, "\n\n".join(text_parts)


def _extract_json(text: str) -> dict:
    try:
        return json.loads(text)
    except Exception:
        pass
    match = re.search(r"```(?:json)?\s*([\s\S]+?)\s*```", text)
    if match:
        try:
            return json.loads(match.group(1))
        except Exception:
            pass
    match = re.search(r"\{[\s\S]+\}", text)
    if match:
        try:
            return json.loads(match.group(0))
        except Exception:
            pass
    return {}


NEWS_TYPE_LABELS = {
    "type_a": "Type A – Nyhetsmelding",
    "type_b": "Type B – Ny-partner",
    "type_c": "Type C – Program/arrangement-rapport",
    "type_d": "Type D – Analytisk artikkel",
    "annet":  "Annet",
}


def generate_questions(
    context: str,
    is_incubator: bool,
    company_name: str,
    news_type: str = "",
    files_data: Optional[list] = None,
) -> dict:
    doc_blocks, files_text = _process_files(files_data or [])

    type_label = NEWS_TYPE_LABELS.get(news_type, "ukjent")
    context_text = f"Valgt artikkeltype: {type_label}\n\nKontekst:\n{context}"
    if files_text:
        context_text += f"\n\nFiler fra selskapet:\n{files_text}"
    if is_incubator and company_name:
        context_text += f"\n\nInkubatorselskap: {company_name}"

    content = doc_blocks + [{"type": "text", "text": context_text}]

    resp = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=1024,
        system="""Du er redaksjonsassistent for Aggrator. Brukeren har ALLEREDE valgt artikkeltype (oppgitt i konteksten). Du skal IKKE klassifisere på nytt.

Din eneste oppgave: identifiser konkrete faktaopplysninger som faktisk mangler for å skrive en god artikkel av den valgte typen.

Returner kun JSON på dette formatet:
{"questions": ["spørsmål 1", "spørsmål 2"]}

Regler:
- Returner FÅ spørsmål. Hvis grunnlaget allerede er godt nok, returner en tom liste: {"questions": []}. Maks 5 spørsmål.
- Spør KUN om konkrete fakta som mangler – aldri om noe som allerede står i konteksten, kulepunktene eller filene.
- Spør ALDRI om bilde, foto eller illustrasjon. Et bilde er allerede lastet opp.
- For Type B (ny-partner): Aggrator ER inkubatoren/vekstpartneren, og selve nyheten er at selskapet nettopp blir en del av Aggrator. Spør ALDRI om selskapet er tilknyttet en inkubator, akselerator eller vekstprogram – svaret er Aggrator.
- Avsender / ansvarlig hos Aggrator er allerede oppgitt. Spør ALDRI om hvem som er ansvarlig, avsender, kontaktperson eller KAM hos Aggrator.
- Still kun innholdsspørsmål om fakta – ikke redaksjonelle/strategiske spørsmål (målgruppe, vinkling, rangering, hva man "håper å oppnå").
- Norsk bokmål""",
        messages=[{"role": "user", "content": content}],
    )

    data = _extract_json(resp.content[0].text)
    return {
        "questions": data.get("questions", []),
    }


def _build_context(news_type: str, fields: dict) -> str:
    f = fields
    if news_type == "type_b":
        base = (
            f"Selskapsnavn: {f['company_name']}\n"
            f"Selskapets hovedkontakt (navn og rolle): {f.get('main_contact') or '[MANGLER: navn og rolle på hovedkontakt]'}\n"
            f"Andre grunnleggere / sentrale teammedlemmer: {f.get('other_team') or '(ingen andre oppgitt – selskapet har kun hovedkontakten over)'}\n"
            f"Ansvarlig hos Aggrator (KAM): {f['responsible']}\n"
            f"Hva driver selskapet med: {f['description']}\n"
            f"Ekstra notater: {f['notes'] or '(ingen)'}"
        )
    elif news_type == "type_a":
        base = (
            f"Hva handler nyheten om: {f['aggrator_what']}\n"
            f"Hvem er involvert: {f['aggrator_who'] or '(ikke oppgitt)'}\n"
            f"Avsender / ansvarlig hos Aggrator: {f.get('responsible') or '(ikke oppgitt)'}\n"
            f"Ekstra notater: {f['notes'] or '(ingen)'}"
        )
    elif news_type == "type_c":
        base = (
            f"Program / arrangement: {f['topic']}\n"
            f"Beskrivelse (fri tekst – la innholdet bestemme struktur, ikke tving dagsoversikt): {f['key_points']}\n"
            f"Avsender / ansvarlig hos Aggrator: {f.get('responsible') or '(ikke oppgitt)'}\n"
            f"Ekstra notater: {f['notes'] or '(ingen)'}"
        )
    elif news_type == "type_d":
        base = (
            f"Tema: {f['topic']}\n"
            f"Nøkkelpunkter: {f['key_points']}\n"
            f"Avsender / ansvarlig hos Aggrator: {f.get('responsible') or '(ikke oppgitt)'}\n"
            f"Ekstra notater: {f['notes'] or '(ingen)'}"
        )
    elif news_type == "annet":
        base = f"Beskrivelse: {f['notes']}"
    else:
        base = ""

    extra = (f.get("extra_context") or "").strip()
    if extra:
        base += f"\n\nUtfyllende svar på oppfølgingsspørsmål:\n{extra}"
    return base


def _build_article_request(
    news_type: str,
    fields: dict,
    corrections: dict,
    files_data: Optional[list],
    image_data: Optional[dict],
) -> tuple[str, list]:
    """Bygger (system, content) for artikkelgenerering. Delt mellom vanlig og
    streamende variant."""
    doc_blocks, files_text = _process_files(files_data or [])

    image_block = None
    if image_data and image_data.get("bytes"):
        image_block = {
            "type": "image",
            "source": {
                "type": "base64",
                "media_type": image_data.get("mime", "image/jpeg"),
                "data": base64.standard_b64encode(image_data["bytes"]).decode("utf-8"),
            },
        }

    lessons_block = ""
    lessons = corrections.get("lessons", [])[-6:]
    if lessons:
        lessons_block = "\n\n## Lærte stilregler fra tidligere korrigeringer:\n"
        for lesson in lessons:
            lessons_block += f"- {lesson['lesson']}\n"

    system = f"""Du er redaksjonsassistent for Aggrator. Skriv en ferdig artikkel på norsk bokmål.

## Retningslinjer:
{GUIDELINES}
{lessons_block}

REGLER:
- Du skal ALLTID levere en komplett, ferdig artikkel med overskrift, ingress, brødtekst, sitatboks og bildetekst. Selv om nesten all informasjon mangler, skriver du artikkelstrukturen og setter inn [MANGLER: ...] der innholdet skal stå. Du skal ALDRI stille spørsmål tilbake til brukeren, ALDRI nekte å skrive, og ALDRI skrive setninger som "Jeg mangler informasjon", "Kan du fylle inn", "Jeg ser at feltene er knappe" eller lignende. Du skal ALDRI lage en liste over hva som mangler. Brukeren kan ikke svare deg – det finnes ikke noe svarfelt. Svaret ditt skal være selve artikkelen, ingenting annet.
- Aldri oppfinn fakta – bruk kun det som er oppgitt
- Bruk [MANGLER: ...] SPARSOMT og kun for opplysninger artikkelen faktisk trenger og som ikke er oppgitt noe sted i konteksten. Marker ALDRI noe som [MANGLER] hvis svaret allerede står i konteksten.
- Kontaktpersonen/KAM hos Aggrator står i konteksten som "Ansvarlig hos Aggrator (KAM)". Bruk det navnet i avslutningssetningen – aldri [MANGLER] der.
- Anta ALDRI at det finnes flere personer eller enheter enn det som er oppgitt. Er det oppgitt én gründer/daglig leder, er det den eneste – etterspør ALDRI "øvrige gründere/eiere/medeiere". Det samme gjelder ansatte, investorer og partnere. Ikke finn på at noe mangler basert på antakelser om hva som kanskje finnes.
- Når du faktisk mangler en nødvendig opplysning: skriv artikkelen videre som normalt og sett inn [MANGLER: kort beskrivelse] inline der opplysningen hører hjemme (f.eks. "Selskapet ble grunnlagt i [MANGLER: stiftelsesår]."). Aldri samle manglene i en egen liste.
- ALDRI lag lenker, URL-er, nettadresser, domener eller e-postadresser. Ikke gjett på et selskaps nettside. Selskaps- og produktnavn skrives som ren tekst (f.eks. GreenRoot AS), aldri som markdown-lenke [navn](url). Bruk kun en URL hvis den eksakte adressen er oppgitt i inndataene.
- Skriv alltid "Vi i Aggrator" (aldri tredjeperson)
- Ingen reklamespråk: "revolusjonerende", "banebrytende", "game-changing" er forbudt
- Sitater i eget innrykk: – [sitat], sier [Navn], [Tittel] i [Selskap].
- Dikt ALDRI opp sitater. Bruk kun sitater som faktisk er oppgitt i inndataene. Er det ikke oppgitt et sitat, skriv [MANGLER: sitat fra Navn] i sitatboksen i stedet for å finne på ett.
- Bildetekst i kursiv til slutt"""

    context = _build_context(news_type, fields)
    user_text = (
        f"Type nyhet: {NEWS_TYPE_LABELS.get(news_type, news_type)}\n\n"
        f"{context}"
    )
    if files_text:
        user_text += f"\n\nFILER FRA SELSKAPET:\n{files_text}"
    if image_block is not None:
        user_text += (
            "\n\nDet følger med ett bilde til artikkelen (vist over). "
            "Skriv en treffende bildetekst i kursiv helt til slutt, basert på hva bildet "
            "faktisk viser sammen med innholdet i artikkelen. Ikke dikt opp detaljer du ikke ser."
        )

    content = doc_blocks
    if image_block is not None:
        content = content + [image_block]
    content = content + [{"type": "text", "text": user_text}]

    return system, content


def generate_article(
    news_type: str,
    fields: dict,
    corrections: dict,
    files_data: Optional[list] = None,
    image_data: Optional[dict] = None,
) -> str:
    system, content = _build_article_request(news_type, fields, corrections, files_data, image_data)
    resp = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=2048,
        system=system,
        messages=[{"role": "user", "content": content}],
    )
    return resp.content[0].text.strip()


def generate_article_stream(
    news_type: str,
    fields: dict,
    corrections: dict,
    files_data: Optional[list] = None,
    image_data: Optional[dict] = None,
):
    """Streamer artikkelteksten bit for bit (generator av tekstbiter).
    Brukes med st.write_stream i appen for synlig progresjon."""
    system, content = _build_article_request(news_type, fields, corrections, files_data, image_data)
    with client.messages.stream(
        model="claude-sonnet-4-6",
        max_tokens=2048,
        system=system,
        messages=[{"role": "user", "content": content}],
    ) as stream:
        for text in stream.text_stream:
            yield text


def extract_lesson(before: str, after: str, corrections_file: Path) -> None:
    resp = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=200,
        system="Du er redaksjonsassistent for Aggrator. Identifiser den viktigste stilregelen som kan læres av denne korrigeringen. Svar med én konkret, handlingsbar regel på norsk bokmål. Ingen forklaring – bare regelen.",
        messages=[{
            "role": "user",
            "content": f"FØR korrektur:\n{before[:1500]}\n\nETTER korrektur:\n{after[:1500]}\n\nHvilken stilregel representerer denne endringen?",
        }],
    )
    lesson = resp.content[0].text.strip()

    data = json.loads(corrections_file.read_text()) if corrections_file.exists() else {"corrections": [], "lessons": []}
    data.setdefault("lessons", []).append({
        "lesson": lesson,
        "ts": datetime.now().isoformat(),
    })
    data["lessons"] = data["lessons"][-20:]
    corrections_file.write_text(json.dumps(data, ensure_ascii=False, indent=2))


_HIDDEN_CHARS = (
    "\u00a0",  # no-break space
    "\u202f",  # narrow no-break space
    "\u2007",  # figure space
    "\u200b",  # zero-width space
    "\u200c",  # zero-width non-joiner
    "\u200d",  # zero-width joiner
    "\ufeff",  # BOM / zero-width no-break space
)


def _clean_credential(value: str) -> str:
    """Fjerner skjulte/ikke-ASCII whitespace-tegn fra credential-verdier
    (typisk kopierings-artefakter), og trimmer kantene."""
    for ch in _HIDDEN_CHARS:
        value = value.replace(ch, "")
    return value.strip()


def send_article_to_vilde(
    article: str,
    responsible: str,
    news_type: str,
    recipient: str = "vilde.wettergreen@aggrator.com",
    image_data: Optional[dict] = None,
) -> tuple[bool, str]:
    """Returns (success, message). Requires SMTP_HOST, SMTP_USER, SMTP_PASSWORD env vars."""
    # Rens bort skjulte tegn som ofte følger med ved kopiering (hardt mellomrom,
    # smalt hardt mellomrom, zero-width, BOM). Disse er aldri tiltenkt i SMTP-
    # credentials, men gir en kryptisk "'ascii' codec can't encode"-feil ved
    # innlogging. App-passord (Gmail/Outlook) vises ofte med slike mellomrom som
    # ikke er del av passordet, så de fjernes.
    smtp_host = _clean_credential(os.environ.get("SMTP_HOST", ""))
    smtp_user = _clean_credential(os.environ.get("SMTP_USER", ""))
    smtp_password = _clean_credential(os.environ.get("SMTP_PASSWORD", ""))

    if not (smtp_host and smtp_user and smtp_password):
        return False, "SMTP ikke konfigurert (sett SMTP_HOST, SMTP_USER og SMTP_PASSWORD)"

    # Fallback: om noe ikke-ASCII fortsatt gjenstår, gi en handlingsrettet melding.
    for label, value in (("SMTP_USER", smtp_user), ("SMTP_PASSWORD", smtp_password)):
        try:
            value.encode("ascii")
        except UnicodeEncodeError:
            return False, (
                f"{label} i .env inneholder et ikke-ASCII-tegn. Skriv verdien på "
                "nytt for hånd i nyhetssaker/interface/.env – ikke kopier-lim den inn."
            )

    smtp_port = int(os.environ.get("SMTP_PORT", "587"))
    type_label = NEWS_TYPE_LABELS.get(news_type, news_type)

    msg = MIMEMultipart("mixed")
    # Emnet inneholder ikke-ASCII (tankestrek, ev. hardt mellomrom) – må utf-8-kodes
    # eksplisitt, ellers feiler send_message med 'ascii' codec can't encode ...
    msg["Subject"] = Header(f"Artikkelutkast til korrektur – {type_label}", "utf-8")
    msg["From"] = smtp_user
    msg["To"] = recipient

    body = (
        f"Hei Vilde,\n\n"
        f"{responsible} har generert et artikkelutkast ({type_label}) som er klart for korrektur.\n\n"
        f"{'—' * 40}\n\n"
        f"{article}\n\n"
        f"{'—' * 40}\n\n"
        f"Bildet til artikkelen ligger som vedlegg.\n\n"
        f"Sendt fra Aggrator Artikkelassistent"
    )
    msg.attach(MIMEText(body, "plain", "utf-8"))

    if image_data and image_data.get("bytes"):
        subtype = image_data.get("mime", "image/jpeg").split("/")[-1]
        img_part = MIMEImage(image_data["bytes"], _subtype=subtype)
        img_part.add_header(
            "Content-Disposition", "attachment",
            filename=image_data.get("name", "bilde.jpg"),
        )
        msg.attach(img_part)

    try:
        with smtplib.SMTP(smtp_host, smtp_port) as server:
            server.ehlo()
            server.starttls()
            server.login(smtp_user, smtp_password)
            server.send_message(msg)
        return True, f"E-post sendt til {recipient}"
    except Exception as e:
        return False, f"Feil ved sending av e-post: {e}"


def update_company_profile(
    company_name: str,
    before: str,
    after: str,
    saker_dir: Path,
) -> None:
    resp = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=512,
        system="""Identifiser faktaopplysninger om selskapet som ble korrigert (navn, produktnavn, titler, tall, datoer osv.).
Returner JSON: {"corrections": [{"field": "beskrivelse av felt", "old": "gammel verdi", "new": "ny verdi"}]}
Hvis ingen faktakorrigeringer: {"corrections": []}""",
        messages=[{
            "role": "user",
            "content": f"FØR selskapsgjennomlesning:\n{before[:1500]}\n\nETTER selskapsgjennomlesning:\n{after[:1500]}",
        }],
    )

    data = _extract_json(resp.content[0].text)
    corrections = data.get("corrections", [])
    if not corrections:
        return

    slug = (
        company_name.lower()
        .replace(" ", "-")
        .replace("æ", "ae")
        .replace("ø", "oe")
        .replace("å", "aa")
    )

    company_dir = saker_dir / slug
    company_dir.mkdir(parents=True, exist_ok=True)
    (company_dir / "kontekst").mkdir(exist_ok=True)
    (company_dir / "output").mkdir(exist_ok=True)

    corrections_text = "\n".join(
        f"- **{c['field']}**: ~~{c['old']}~~ → {c['new']}" for c in corrections
    )
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M")
    update_block = f"\n\n## Faktakorrigeringer fra selskap ({timestamp})\n\n{corrections_text}\n"

    claude_md = company_dir / "CLAUDE.md"
    if claude_md.exists():
        claude_md.write_text(claude_md.read_text() + update_block)
    else:
        claude_md.write_text(f"# {company_name}\n{update_block}")
